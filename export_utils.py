"""
분석 결과를 Word(.docx) / PPT(.pptx) 파일(바이트)로 변환하는 모듈.
app.py의 result(dict)를 입력받는다.
result = {name, overview(dict), fin(dict), disclosures(list), report(markdown str), ...}
"""

import io
import re
from datetime import datetime

from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PRGB

MARKET = {"Y": "코스피", "K": "코스닥", "N": "코넥스", "E": "기타"}


def _clean(text):
    """마크다운 강조 기호 제거."""
    return text.replace("**", "").replace("`", "").strip()


def _fmt_eok(v):
    if v is None:
        return "-"
    try:
        return f"{v / 1e8:,.1f}"
    except Exception:
        return str(v)


def _parse_sections(md):
    """마크다운 리포트를 (섹션제목, [줄들]) 목록으로 파싱."""
    sections = []
    title, lines = None, []
    for raw in (md or "").splitlines():
        line = raw.rstrip()
        if line.startswith("## "):
            if title is not None:
                sections.append((title, lines))
            title, lines = _clean(line[3:]), []
        elif line.startswith("# "):
            continue
        elif line.strip():
            lines.append(_clean(line.lstrip("-• ").strip()))
    if title is not None:
        sections.append((title, lines))
    return sections


# --------------------------------------------------------------------- WORD
def build_docx(res):
    doc = Document()

    title = doc.add_heading(f"{res.get('name','')} 분석 리포트", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph(
        f"데이터 출처: 금융감독원 OpenDART · 생성일 {datetime.now():%Y-%m-%d}"
    )
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for r in sub.runs:
        r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    ov = res.get("overview") or {}
    if ov:
        doc.add_heading("기업 개황", level=1)
        info = [
            ("대표자", ov.get("ceo_nm", "-")),
            ("설립일", ov.get("est_dt", "-")),
            ("시장", MARKET.get(ov.get("corp_cls", ""), ov.get("corp_cls", "-"))),
            ("종목코드", ov.get("stock_code") or "비상장"),
            ("주소", ov.get("adres", "-")),
        ]
        for k, v in info:
            p = doc.add_paragraph()
            p.add_run(f"{k}: ").bold = True
            p.add_run(str(v))

    fin = res.get("fin")
    if fin:
        doc.add_heading("재무제표 (최근 3개년, 단위: 억원)", level=1)
        years = fin["years"]
        table = doc.add_table(rows=1, cols=1 + len(years))
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        hdr[0].text = "계정"
        for i, y in enumerate(years):
            hdr[i + 1].text = f"{y}년"
        for name, vals in fin["rows"].items():
            cells = table.add_row().cells
            cells[0].text = name
            for i, v in enumerate(vals):
                cells[i + 1].text = _fmt_eok(v)
        note = doc.add_paragraph(
            f"기준: {'연결재무제표' if fin.get('fs_div') == 'CFS' else '별도재무제표'}"
        )
        for r in note.runs:
            r.font.size = Pt(9)
            r.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    doc.add_heading("AI 종합 분석", level=1)
    for stitle, slines in _parse_sections(res.get("report", "")):
        doc.add_heading(stitle, level=2)
        for ln in slines:
            doc.add_paragraph(ln, style="List Bullet")

    disclosures = res.get("disclosures") or []
    if disclosures:
        doc.add_heading("최근 공시 목록", level=1)
        for d in disclosures[:20]:
            doc.add_paragraph(
                f"{d.get('rcept_dt','')} · {d.get('report_nm','')}",
                style="List Bullet",
            )

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------- PPT
def _add_bullets(placeholder, lines):
    tf = placeholder.text_frame
    tf.word_wrap = True
    if not lines:
        tf.text = "자료에 정보가 부족함"
        return
    tf.text = lines[0]
    for ln in lines[1:]:
        p = tf.add_paragraph()
        p.text = ln


def build_pptx(res):
    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # 표지
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = f"{res.get('name','')}\nDART 공시 분석 리포트"
    ov = res.get("overview") or {}
    subtitle = (
        f"대표자 {ov.get('ceo_nm','-')} · "
        f"{MARKET.get(ov.get('corp_cls',''), ov.get('corp_cls','-'))} · "
        f"생성일 {datetime.now():%Y-%m-%d}"
    )
    if len(s.placeholders) > 1:
        s.placeholders[1].text = subtitle

    # 재무 요약 표
    fin = res.get("fin")
    if fin:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        tb.text_frame.text = "재무제표 (최근 3개년, 단위: 억원)"
        tb.text_frame.paragraphs[0].font.size = PPt(24)
        tb.text_frame.paragraphs[0].font.bold = True

        years = fin["years"]
        rows = len(fin["rows"]) + 1
        cols = 1 + len(years)
        gt = s.shapes.add_table(
            rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5 * rows)
        ).table
        gt.cell(0, 0).text = "계정"
        for i, y in enumerate(years):
            gt.cell(0, i + 1).text = f"{y}년"
        for ri, (name, vals) in enumerate(fin["rows"].items(), start=1):
            gt.cell(ri, 0).text = name
            for i, v in enumerate(vals):
                gt.cell(ri, i + 1).text = _fmt_eok(v)
        for r in range(rows):
            for c in range(cols):
                for p in gt.cell(r, c).text_frame.paragraphs:
                    p.font.size = PPt(12)

    # 분석 섹션별 슬라이드
    for stitle, slines in _parse_sections(res.get("report", "")):
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = stitle
        _add_bullets(s.placeholders[1], slines[:12])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
